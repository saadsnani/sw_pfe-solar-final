from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Couleurs
dark_bg = RGBColor(20, 24, 28)
emerald = RGBColor(0, 201, 167)
white = RGBColor(255, 255, 255)

def set_dark_theme(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = white
            paragraph.alignment = PP_ALIGN.LEFT
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = dark_bg

def add_image_placeholder(slide, left=Inches(6), top=Inches(1.5), width=Inches(3), height=Inches(3)):
    shape = slide.shapes.add_shape(
        1, left, top, width, height  # 1 = Rectangle
    )
    shape.text = "IMAGE_PLACEHOLDER"
    shape.text_frame.paragraphs[0].font.size = Pt(14)
    shape.text_frame.paragraphs[0].font.color.rgb = emerald
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(30, 35, 40)
    shape.line.color.rgb = emerald

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slides_content = [
    # Section 1: Introduction & Contexte
    ("Smart EMS : Monitoring Photovoltaïque Intelligent via IA", "Projet de fin d’études\nSaad Snani & Walid El Halouat\n\nThème : Monitoring intelligent, énergie propre, IA.\n", True),
    ("Problématique du Solaire Off-grid", "Défis :\n- Autonomie énergétique\n- Fiabilité du monitoring\n- Maintenance préventive\n", True),
    ("Objectifs du Projet", "• Monitoring temps réel\n• Optimisation énergétique\n• Sécurité et maintenance intelligente\n", True),

    # Section 2: Puissance & Sécurité
    ("Dimensionnement du Système", "• Panneau solaire : 130W\n• Batterie : 12V\n• Régulateur MPPT\n• Calculs adaptés à l’autonomie Off-grid\n", True),
    ("Calculs de Puissance", "Formules utilisées :\n- Puissance = Tension x Courant\n- Capacité batterie = (Conso x Autonomie) / Rendement\n- Prise en compte des pertes\n", True),
    ("Régulateur MPPT (98%)", "• Algorithme Perturb & Observe (P&O)\n• Suivi du point de puissance maximale\n• Efficacité mesurée : 98%\n", True),
    ("Normes de Protection", "• Fusibles DC adaptés\n• Disjoncteurs pour sécurité\n• Respect des normes électriques solaires\n", True),
    ("Schéma de Protection", "Schéma Proteus des protections électriques.\n", True),

    # Section 3: Hardware & IoT
    ("Architecture Dual-Core", "• Arduino Mega : Acquisition capteurs\n• ESP32 : Communication WiFi & traitement\n• Communication série UART entre les deux\n", True),
    ("Communication Série Mega/ESP32", "• Protocole UART\n• Fiabilité des échanges de données\n• Synchronisation acquisition/émission\n", True),
    ("Capteur de Courant ACS712", "• Mesure précise du courant\n• Intégration sur bus DC\n• Calibration et filtrage logiciel\n", True),
    ("Capteur Température DS18B20", "• Bus OneWire\n• Robustesse et précision\n• Surveillance thermique batterie/panneau\n", True),
    ("Schéma d’Acquisition", "Schéma Proteus de l’acquisition capteurs.\n", True),

    # Section 4: Software Stack
    ("Dashboard Web (Next.js 14)", "• Frontend : Next.js 14, Tailwind CSS\n• Visualisation temps réel\n• Mode sombre, responsive\n", True),
    ("Base de Données MongoDB", "• Stockage sécurisé des mesures\n• Requêtes optimisées\n• Historique et analyse des données\n", True),
    ("Architecture Logicielle", "• Flux : Capteurs → ESP32 → API → Dashboard\n• API REST sécurisée\n• Gestion des utilisateurs\n", True),
    ("Nom de Domaine .me", "• Hébergement sur domaine personnalisé\n• Accès sécurisé et professionnel\n", True),

    # Section 5: Innovation IA (Gemini Pro)
    ("Diagnostic en Langage Naturel (Gemini Pro)", "• IA générative intégrée au dashboard\n• Explications claires des anomalies\n• Interaction utilisateur simplifiée\n", True),
    ("Maintenance Prédictive", "• Détection automatique des dérives\n• Alertes intelligentes\n• Réduction des pannes\n", True),
    ("Modes Intelligents (Eco/Boost)", "• Mode Eco : optimisation consommation\n• Mode Boost : performance maximale\n• Sélection automatique selon contexte\n", True),
    ("Intégration IA dans le Dashboard", "• Interface dédiée aux diagnostics IA\n• Visualisation des recommandations\n", True),

    # Section 6: Résultats & Conclusion
    ("Validation Expérimentale", "• Tests sur banc réel\n• Mesures de performance et fiabilité\n• Comparaison avec solutions classiques\n", True),
    ("Impact Environnemental", "• Réduction des émissions CO2\n• Promotion de l’énergie propre\n• Contribution à la transition énergétique\n", True),
    ("Perspectives", "• Conception d’un PCB dédié\n• Améliorations futures (Edge AI, Cloud)\n", True),
    ("Remerciements & Questions", "Merci pour votre attention !\nContact : saad.snani@exemple.com\n", False),
]

for title, content, has_placeholder in slides_content:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = emerald

    # Content
    content_shape = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.5), Inches(4.5))
    content_frame = content_shape.text_frame
    content_frame.text = content
    content_frame.paragraphs[0].font.size = Pt(22)
    content_frame.paragraphs[0].font.color.rgb = white

    # Placeholder for image
    if has_placeholder:
        add_image_placeholder(slide)

    set_dark_theme(slide)

prs.save("Smart-EMS-PFE-Presentation.pptx")
print("Présentation générée : Smart-EMS-PFE-Presentation.pptx")
